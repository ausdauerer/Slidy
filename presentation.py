from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
from pptx.dml.color import RGBColor


def add_title_presentation_slide(prs, slide):
    title = slide.get("title", {
        "text": "Title",
        "font": "Arial",
        "color": (0, 0, 0)
    })
    subtitle = slide.get("subtitle", None)

    slide_layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(slide_layout)

    title_shape = s.shapes.title
    title_shape.text = title["text"]

    if subtitle:
        subtitle_shape = s.placeholders[1]
        subtitle_shape.text = subtitle["text"]
        if subtitle_shape.text_frame.paragraphs:
            run = subtitle_shape.text_frame.paragraphs[0].runs[0]
            run.font.name = subtitle["font"]
            run.font.bold = False
            run.font.size = Pt(28)
            run.font.color.rgb = RGBColor(subtitle["color"][0], subtitle["color"][1], subtitle["color"][2])

    if title_shape.text_frame.paragraphs:
        run = title_shape.text_frame.paragraphs[0].runs[0]
        run.font.name = title["font"]
        run.font.bold = True
        run.font.size = Pt(44)
        run.font.color.rgb = RGBColor(title["color"][0], title["color"][1], title["color"][2])

def add_bullet_presentation_slide(prs, slide):
    title = slide.get("title", {})
    bullets = slide.get("bullet_points", [] )

    slide_layout = prs.slide_layouts[1]
    s = prs.slides.add_slide(slide_layout)

    title_shape = s.shapes.title
    title_shape.text = title["text"]

    if title_shape.text_frame.paragraphs:
        run = title_shape.text_frame.paragraphs[0].runs[0]
        run.font.name = title["font"]
        run.font.bold = False
        run.font.color.rgb = RGBColor(title["color"][0], title["color"][1], title["color"][2])

    body = s.placeholders[1].text_frame
    for bullet in bullets:
        p = body.add_paragraph()
        p.text = bullet["text"]
        p.level = 0
        
        if p.runs:
            run = p.runs[0]
            run.font.name = bullet["font"]
            run.font.size = Pt(24)
            run.font.bold = False
            run.font.color.rgb = RGBColor(bullet["color"][0], bullet["color"][1], bullet["color"][2])

def add_image_with_context_presentation_slide(prs, slide):
    title = slide.get("title", {})
    images = slide.get("images", [])
    content = slide.get("content", {})

    slide_layout = prs.slide_layouts[5]
    s = prs.slides.add_slide(slide_layout)

    title_shape = s.shapes.title
    title_shape.text = title["text"]

    if title_shape.text_frame.paragraphs:
        run = title_shape.text_frame.paragraphs[0].runs[0]
        run.font.name = title["font"]
        run.font.bold = False
        run.font.color.rgb = RGBColor(title["color"][0], title["color"][1], title["color"][2])

    margin_left = Inches(0.5)
    spacing = Inches(0.3)
    max_width = prs.slide_width - 2 * margin_left
    current_top = title_shape.top + title_shape.height + spacing

    if images:
        img_path = images[0]

        # Get image size in pixels
        from PIL import Image

        with Image.open(img_path) as img:
            img_width_px, img_height_px = img.size
            img_ratio = img_width_px / img_height_px

            img_left = (prs.slide_width - 0) / 2

            # Scale image to fit max_width without exceeding slide height
            target_width = max_width
            target_height = target_width / img_ratio

            if current_top + target_height > prs.slide_height:  
                # If image too tall, scale down by height instead
                target_height = prs.slide_height - current_top - Inches(2)
                target_width = target_height * img_ratio

            # Place the image
            s.shapes.add_picture(img_path, 0, current_top, width=target_width, height=target_height)
            current_top += target_height + spacing

    # Add Content
    if content:
        textbox_height = prs.slide_height - current_top - Inches(0.5)
        textbox = s.shapes.add_textbox(margin_left, current_top, max_width, textbox_height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.text = content["text"]

        # Set custom font and color if specified
        font_name = content.get("font", "Arial")
        font_color = content.get("color", (0, 0, 0))
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = content["font"]
                run.font.color.rgb = RGBColor(content["color"][0], content["color"][1], content["color"][2])

def add_two_column_presentation_slide(prs, slide):
    title = slide.get("title", {})
    content=slide.get("content", {})
    images=slide.get("images", [])

    slide_layout = prs.slide_layouts[5] 
    s = prs.slides.add_slide(slide_layout)

    title_shape = s.shapes.title
    title_shape.text = title["text"]

    if title_shape.text_frame.paragraphs:
        run = title_shape.text_frame.paragraphs[0].runs[0]
        run.font.name = title["font"]
        run.font.bold = False
        run.font.color.rgb = RGBColor(title["color"][0], title["color"][1], title["color"][2])

    # Define margins and column widths
    margin_left = Inches(0.5)
    margin_right = Inches(0.5)
    spacing = Inches(0.3)
    col_width = (prs.slide_width - margin_left - margin_right - spacing) / 2
    col_height = prs.slide_height - title_shape.top - title_shape.height - Inches(1)
    top = title_shape.top + title_shape.height + spacing

    if content:
        textbox = s.shapes.add_textbox(margin_left, top, col_width, col_height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.text = content["text"]

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = content["font"]
                run.font.color.rgb = RGBColor(content["color"][0], content["color"][1], content["color"][2])

    # Right column: image (if any)
    if images:
        img_path = images[0]
        with Image.open(img_path) as img:
            img_width_px, img_height_px = img.size

            img_ratio = img_width_px / img_height_px

            # Fit image to column width, preserve aspect ratio
            target_width = col_width
            target_height = target_width / img_ratio
            if target_height > col_height:
                target_height = col_height
                target_width = target_height * img_ratio

            image_left = margin_left + col_width + spacing
            s.shapes.add_picture(img_path, image_left, top, width=target_width, height=target_height)
    

def create_presentation(slide_data, output_file="presentation.pptx"):
    prs = Presentation()

    for slide in slide_data["slides"]:
        template = slide.get("template", "Title Slide")

        if template == "Title Slide":
            add_title_presentation_slide(prs, slide)

        elif template == "Bullet Points":
            add_bullet_presentation_slide(prs, slide)

        elif template == "Image with Content":
            add_image_with_context_presentation_slide(prs, slide)

        elif template == "2 Column Layout":
            add_two_column_presentation_slide(prs, slide)
        else:
            # default to title slide if unknown template
            s = prs.slides.add_slide(prs.slide_layouts[0])
            s.shapes.title.text = title_text

    prs.save("presentations/" + output_file+".pptx")

    return "presentations/" + output_file+".pptx"

